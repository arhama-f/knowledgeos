import io

from pptx import Presentation

from app.services.extractors.base import ExtractedPage, Extractor


class PptxExtractor(Extractor):
    def extract(self, content: bytes) -> list[ExtractedPage]:
        presentation = Presentation(io.BytesIO(content))
        pages = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            texts = [
                shape.text_frame.text.strip()
                for shape in slide.shapes
                if shape.has_text_frame and shape.text_frame.text.strip()
            ]
            if texts:
                pages.append(ExtractedPage(page_number=slide_index, text="\n\n".join(texts)))
        return pages
